from typing import Optional, Dict, Any, List, Tuple
from pathlib import Path
import pandas as pd
import threading
from queue import Queue
from pptx import Presentation
from pptx.dml.color import RGBColor
import shutil
from datetime import datetime
import os
from .terminology import TerminologyManager
from .prompt_manager import PromptManager
from concurrent.futures import ThreadPoolExecutor, as_completed
import time
from pptx.util import Pt
import re
from langdetect import detect
import string
import math

class Translator:
    def __init__(self, config):
        self.config = config
        self.terminology_manager = TerminologyManager()
        self.prompt_manager = PromptManager()
        self.progress_callback = None
        self.status_callback = None
        self.excel_queue = Queue()
        self.translation_lock = threading.Lock()
        self.max_retries = 3  # 最大重试次数
        self.executor = ThreadPoolExecutor(max_workers=3)  # 设置线程池大小
        self.init_client()
        self.translation_cache = {}  # 添加翻译缓存
        self.load_terminology_rules()  # 加载术语规则
        
    def update_config(self, new_config: Dict[str, Any]):
        """更新配置并重新初始化客户端"""
        try:
            # 更新配置
            self.config.config.update(new_config)
            
            # 重新初始化客户端
            self.init_client()
        except Exception as e:
            print(f"更新配置失败: {str(e)}")
            raise

    def init_client(self):
        """初始化翻译客户端"""
        try:
            if self.config.config.get("use_online", True):
                from ..api.zhipuai_client import ZhipuAIClient
                self.client = ZhipuAIClient(self.config.config["api_key"])
            else:
                from ..api.ollama_client import OllamaClient
                self.client = OllamaClient(self.config.config["server_url"])
        except Exception as e:
            print(f"初始化翻译客户端失败: {str(e)}")
            raise

    def translate_ppt(self, input_path: str, use_terminology: bool = False) -> Dict[str, Any]:
        try:
            # 获取文件名和时间戳
            file_name = Path(input_path).stem
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            
            # 创建输出文件路径
            output_dir = Path("src/output")
            output_dir.mkdir(parents=True, exist_ok=True)
            
            output_path = output_dir / f"{file_name}_{timestamp}_translated.pptx"
            excel_path = output_dir / f"{file_name}_{timestamp}_record.xlsx"
            
            # 创建空的Excel文件
            pd.DataFrame(columns=["页码", "位置", "格式", "原文", "翻译结果"]).to_excel(excel_path, index=False)
            
            # 加载PPT
            prs = Presentation(input_path)
            total_slides = len(prs.slides)
            print(f"\n开始翻译PPT文件: {file_name}")
            print(f"总页数: {total_slides}")
            print("="*50)
            
            # 添加进度回调
            if self.progress_callback:
                self.progress_callback(f"\n开始翻译PPT文件: {file_name}")
                self.progress_callback(f"总页数: {total_slides}")
                self.progress_callback("="*50)
            
            # 启动Excel写入线程
            self.excel_thread = threading.Thread(
                target=self._excel_writer_thread,
                args=(excel_path,),
                daemon=True
            )
            self.excel_thread.start()
            
            # 计算总任务数
            total_tasks = self._calculate_total_tasks(prs)
            completed_tasks = 0
            
            # 创建进度追踪变量
            current_progress = {
                "slide": 0,
                "shape": 0,
                "text": "",
                "translation": ""
            }
            
            # 处理每一页
            for slide_idx, slide in enumerate(prs.slides, 1):
                print(f"\n正在处理第 {slide_idx}/{total_slides} 页")
                # 添加进度回调
                if self.progress_callback:
                    self.progress_callback(f"\n正在处理第 {slide_idx}/{total_slides} 页")
                
                current_progress["slide"] = slide_idx
                
                for shape_idx, shape in enumerate(slide.shapes, 1):
                    current_progress["shape"] = shape_idx
                    
                    # 处理形状并等待完成
                    result = self._process_shape_with_progress(
                        shape, 
                        slide_idx, 
                        use_terminology,
                        current_progress
                    )
                    
                    if result:
                        # 更新进度信息
                        self._update_progress(current_progress)
                        
                        # 保存当前进度
                        try:
                            prs.save(output_path)
                        except Exception as e:
                            print(f"保存进度时出错: {str(e)}")
            
                print(f"第 {slide_idx} 页处理完成")
                print("-"*50)
            
            # 等待所有Excel写入任务完成
            self.excel_queue.put(None)  # 发送结束信号
            self.excel_queue.join()
            self.excel_thread.join()
            
            print("\n翻译完成！")
            print(f"输出文件: {output_path}")
            print(f"翻译记录: {excel_path}")
            print("="*50)
            
            # 完成翻译后显示复检完成消息
            final_status = {
                "current_location": "翻译和复检已完成",
                "original_text": "",
                "translated_text": "请查收翻译结果文件",
                "preview_data": pd.read_excel(excel_path).values.tolist(),
                "output_file": str(output_path)
            }
            
            if self.status_callback:
                self.status_callback(final_status)
                
            return final_status
            
        except Exception as e:
            print(f"\n翻译失败: {str(e)}")
            raise Exception(f"PPT翻译失败: {str(e)}")
    
    def _calculate_total_tasks(self, prs) -> int:
        """计算总任务数"""
        total = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    total += len(table.rows) * len(table.columns)
                if hasattr(shape, "text_frame"):
                    total += 1
        return total
    
    def _process_shape_with_progress(self, shape, slide_idx: int, use_terminology: bool, progress: dict) -> bool:
        """处理单个形状并更新进度"""
        try:
            texts_to_translate = []
            text_locations = []

            # 收集需要翻译的文本
            if shape.has_table:
                print(f"\n处理表格...")
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        if cell.text.strip():
                            location = f"Slide{slide_idx}-Table-Row{row_idx+1}-Column{col_idx+1}"
                            texts_to_translate.append(cell.text.strip())
                            text_locations.append((cell.text_frame, location, shape))
                            print(f"发现表格文本 - 位置: {location}")

            elif hasattr(shape, "text_frame") and shape.text_frame.text.strip():
                location = f"Slide{slide_idx}-TextBox"
                texts_to_translate.append(shape.text_frame.text.strip())
                text_locations.append((shape.text_frame, location, shape))
                print(f"发现文本框 - 位置: {location}")

            # 批量翻译文本
            if texts_to_translate:
                translations = self.batch_translate_texts(texts_to_translate, use_terminology)

                # 更新文本框
                for (text_frame, location, shape), translated_text in zip(text_locations, translations):
                    self._update_text_frame_with_translation(
                        text_frame, location, shape, translated_text, progress
                    )

            # 处理组合形状
            if hasattr(shape, "group_items"):
                print(f"\n处理组合形状...")
                for item in shape.group_items:
                    self._process_shape_with_progress(item, slide_idx, use_terminology, progress)

            return True

        except Exception as e:
            print(f"处理形状时出错: {str(e)}")
            return False

    def _update_text_frame_with_translation(self, text_frame, location: str, shape, translated_text: str, progress: dict):
        """更新文本框内容并记录翻译结果"""
        try:
            original_text = text_frame.text.strip()
            progress["text"] = original_text
            progress["translation"] = translated_text

            # 更新文本框
            self._adjust_text_frame(text_frame, translated_text)

            # 添加到Excel队列
            self.excel_queue.put([{
                "页码": location.split("-")[0].replace("幻灯片", ""),
                "位置": location,
                "格式": "文本框" if not hasattr(shape, "table") else "表格",
                "原文": original_text,
                "翻译结果": translated_text
            }])

        except Exception as e:
            print(f"更新文本框时出错: {str(e)}")
    
    def _update_progress(self, progress: dict):
        """更新进度信息到UI"""
        if self.status_callback:
            status = {
                "current_location": f"幻灯片 {progress['slide']} - 形状 {progress['shape']}",
                "original_text": progress["text"],
                "translated_text": progress["translation"]
            }
            self.status_callback(status)
    
    def _process_and_translate_text(self, text: str, location: str, use_terminology: bool) -> Dict[str, str]:
        """处理和翻译文本"""
        # 如果文本只包含数字或英文字母，直接返回原文
        if not self._is_translatable(text):
            return {
                "location": location,
                "original": text,
                "first_translation": text,
                "verified_translation": text,
                "final_translation": text
            }
        
        # 应用术语库（如果启用）
        if use_terminology:
            text = self.terminology_manager.apply_terminology(text)
        
        # 进行翻译
        first_translation = self.translate_text(text)
        verified_translation = self.translate_text(first_translation)  # 复检翻译
        
        return {
            "location": location,
            "original": text,
            "first_translation": first_translation,
            "verified_translation": verified_translation,
            "final_translation": verified_translation
        }
    
    def _excel_writer_thread(self, excel_path: Path):
        """Excel写入线程"""
        records = []
        while True:
            batch = self.excel_queue.get()
            if batch is None:  # 结束信号
                break
            records.extend(batch)
            self.excel_queue.task_done()
        
        # 创建DataFrame并保存
        df = pd.DataFrame(records, columns=[
            "页码", "位置", "格式", "原文", "翻译结果"
        ])
        df.to_excel(excel_path, index=False)
        self.excel_queue.task_done()
    
    def _is_translatable(self, text: str) -> bool:
        """检查文本是否需要翻译"""
        # 如果文本是纯数字或英文编号，则不需要翻译
        if text.strip().isdigit():  # 纯数字
            return False
        
        # 检查是否为英文编号格式（如 A01, B-02 等）
        if re.match(r'^[A-Za-z0-9\-_.]+$', text.strip()):
            return False
        
        # 如果目标是英文，检查是否包含中文字符
        if self.config.config["target_language"] == "英文":
            return self.prompt_manager.has_chinese(text)
        
        # 如果目标是中文，检查是否包含英文字符
        return bool(re.search(r'[a-zA-Z]', text))
    
    def _adjust_text_frame(self, text_frame, new_text: str):
        """调整文本框内容，保持原有格式"""
        try:
            # 更新文本
            text_frame.text = new_text
            
            # 遍历所有段落和文字块，统一设置格式
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    # 设置字体大小为7.0磅
                    if run.font.size and run.font.size.pt > 7.0:
                        run.font.size = Pt(7.0)
                    # 如果未设置字体大小，则设置为7.0磅
                    elif not run.font.size:
                        run.font.size = Pt(7.0)
                        
                    # 统一设置字体颜色为黑色
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    
        except Exception as e:
            print(f"调整文本框格式时出错: {str(e)}")

    def translate_text(self, text: str, use_terminology: bool = False) -> str:
        """改进的翻译方法"""
        try:
            # 1. 检查缓存
            cache_key = f"{text}_{use_terminology}"
            if cache_key in self.translation_cache:
                return self.translation_cache[cache_key]

            # 2. 预处理文本
            processed_text = self._preprocess_text(text)
            
            # 3. 获取翻译
            translation = self._get_translation(processed_text, use_terminology)
            
            # 4. 后处理
            translation = self._postprocess_translation(translation, text)
            
            # 5. 缓存结果
            self.translation_cache[cache_key] = translation
            
            return translation
            
        except Exception as e:
            print(f"翻译失败: {str(e)}")
            return text

    def _preprocess_text(self, text: str) -> str:
        """文本预处理"""
        # 1. 保护特殊格式
        text = self._protect_special_format(text)
        
        # 2. 标准化空格和标点
        text = self._standardize_format(text)
        
        # 3. 处理特殊标记
        text = self._process_special_markers(text)
        
        return text

    def _protect_special_format(self, text: str) -> str:
        """保护特殊格式"""
        # 保护数字和单位
        text = re.sub(r'(\d+)([a-zA-Z°℃]+)', r'__NUM__\1__UNIT__\2__END__', text)
        
        # 保护括号内容
        text = re.sub(r'（([^）]+)）', r'__LEFT__\1__RIGHT__', text)
        
        return text

    def _standardize_format(self, text: str) -> str:
        """标准化格式"""
        # 统一中文标点到英文
        punctuation_map = {
            '，': ',',
            '。': '.',
            '：': ':',
            '；': ';'
        }
        for cn, en in punctuation_map.items():
            text = text.replace(cn, en)
            
        return text

    def _postprocess_translation(self, translation: str, original: str) -> str:
        """翻译后处理"""
        # 1. 恢复特殊格式
        translation = self._restore_special_format(translation)
        
        # 2. 应用术语规则
        translation = self._apply_terminology_rules(translation)
        
        # 3. 格式化输出
        translation = self._format_output(translation)
        
        # 4. 验证翻译质量
        if not self._validate_translation(translation, original):
            translation = self._retry_translation(original)
            
        return translation

    def _apply_terminology_rules(self, text: str) -> str:
        """应用术语规则"""
        # 应用强制替换规则
        for old, new in self.terminology_rules["force_replace"].items():
            text = re.sub(rf'\b{old}\b', new, text, flags=re.IGNORECASE)
            
        # 应用上下文相关规则
        for term, rules in self.terminology_rules["context_replace"].items():
            if term in text:
                for context, replacement in rules["contexts"].items():
                    if context in text:
                        text = text.replace(context, replacement)
                        break
                else:
                    text = text.replace(term, rules["default"])
                    
        return text

    def _format_output(self, text: str) -> str:
        """格式化输出"""
        # 1. 标准化空格
        text = re.sub(r'\s+', ' ', text)
        
        # 2. 修正数字和单位之间的空格
        text = re.sub(r'(\d+)\s*([a-zA-Z°℃Ω]+)', r'\1 \2', text)
        
        # 3. 确保正确的结束标记
        if '[START]' in text and not text.endswith('[END]'):
            text = text.replace('[/START]', '[END]')
            if not text.endswith('[END]'):
                text += '[END]'
                
        return text.strip()

    def translate_text_with_retry(self, text: str, use_terminology: bool = False, retry_count: int = 0) -> str:
        """带重试机制的文本翻译"""
        if not text.strip() or not self._is_translatable(text):
            return text
        
        translations = []  # 存储多次翻译结果
        max_retries = 3
        
        # 进行多次翻译尝试
        while retry_count < max_retries and len(translations) < 3:
            try:
                # 使用不同的提示词变体
                messages = self.prompt_manager.get_variant_prompt(text, retry_count)
                translated = self.client.translate(messages, self.config.config["model"])
                
                # 验证翻译结果
                if translated and translated.strip():
                    if self.config.config["target_language"] == "英文":
                        # 对于英文翻译，确保没有中文
                        if not self.prompt_manager.has_chinese(translated):
                            translations.append(translated)
                    else:
                        # 对于其他语言的翻译，直接添加
                        translations.append(translated)
                
                retry_count += 1
                time.sleep(1)  # 添加短暂延迟避免请求过快
                
            except Exception as e:
                print(f"翻译重试 {retry_count + 1} 失败: {str(e)}")
                retry_count += 1
                time.sleep(2)
        
        # 如果有任何有效的翻译结果，进行选择
        if translations:
            # 使用复检提示词评估翻译质量
            best_translation = self._select_best_translation(text, translations)
            return best_translation
            
        return text  # 如果所有尝试都失败，返回原文

    def _select_best_translation(self, original_text: str, translations: list) -> str:
        """选择最佳翻译结果"""
        try:
            # 如果只有一个翻译结果，直接返回
            if len(translations) == 1:
                return translations[0]
            
            # 如果有两个或更多翻译结果，构建评估提示词
            evaluation_prompt = f"""Please evaluate the quality of the following translation versions and select the most accurate and fluent one.
Original text: {original_text}

"""
            # 动态添加翻译版本
            for i, trans in enumerate(translations, 1):
                evaluation_prompt += f"Version {i}: {trans}\n"
            
            evaluation_prompt += "\nPlease return the best translation directly without any explanation."

            messages = [
                {"role": "system", "content": "You are a professional translation evaluator. Please select the best translation version."},
                {"role": "user", "content": evaluation_prompt}
            ]
            
            # 获取评估结果
            best_translation = self.client.translate(messages, self.config.config["model"])
            
            # 如果评估结果不在translations中，返回最后一个翻译
            if best_translation not in translations:
                return translations[-1]
            
            return best_translation
            
        except Exception as e:
            print(f"选择最佳翻译失败: {str(e)}")
            # 出错时返回最后一个翻译结果（通常是最好的一次）
            return translations[-1] if translations else original_text

    def translate_text_with_verification(self, text: str, use_terminology: bool = False) -> str:
        """增强的多轮校验翻译"""
        try:
            if not text.strip() or not self._is_translatable(text):
                return text
            
            # 1. 第一轮翻译
            first_translation = self.translate_text(text, use_terminology)
            
            # 移除 [START] 和 [END] 标记
            first_translation = re.sub(r'\[START\]|\[END\]', '', first_translation).strip()
            
            # 2. 验证翻译结果
            is_valid, reason = self._is_valid_translation_with_reason(text, first_translation)
            
            if not is_valid:
                # 使用调整后的策略重新翻译
                adjusted_text = self._adjust_translation_strategy(text, [reason])
                alternative_translation = self.translate_text(adjusted_text, use_terminology)
                alternative_translation = re.sub(r'\[START\]|\[END\]', '', alternative_translation).strip()
                
                # 再次验证
                is_valid_alt, _ = self._is_valid_translation_with_reason(text, alternative_translation)
                
                # 选择更好的翻译结果
                if is_valid_alt:
                    return alternative_translation
                
                # 如果两次翻译都失败，尝试分词翻译
                return self._get_best_candidate(text, [reason])
                
            return first_translation
            
        except Exception as e:
            print(f"翻译验证失败: {str(e)}")
            return text

    def _verify_terminology(self, translation: str) -> bool:
        """验证术语准确性"""
        try:
            if not hasattr(self, 'terminology_manager') or not self.terminology_manager.terminology:
                return True
            
            # 检查所有术语是否正确翻译
            for cn_term, en_term in self.terminology_manager.terminology.items():
                if cn_term in translation and en_term not in translation:
                    return False
                
            return True
        except:
            return True

    def _verify_fluency(self, translation: str) -> bool:
        """验证语言流畅度"""
        try:
            # 1. 检查是否包含基本的语言结构
            words = translation.split()
            if len(words) < 2:
                return False
            
            # 2. 检查标点符号使用
            if translation.count('.') > 1 and not any(char in translation for char in ',;:'):
                return False
            
            # 3. 检查重复词
            word_counts = {}
            for word in words:
                word_counts[word] = word_counts.get(word, 0) + 1
                if word_counts[word] > 3:  # 同一个词出现超过3次
                    return False
            
            return True
        except:
            return True

    def _verify_context(self, translation: str) -> bool:
        """验证上下文一致性"""
        try:
            # 简单的上下文检查
            if translation.endswith(('.', ':', ';')):
                return True
            
            # 检查括号匹配
            if translation.count('(') != translation.count(')'):
                return False
            
            return True
        except:
            return True

    def _score_translation(self, original: str, translation: str) -> float:
        try:
            score = 0.0
            
            # 1. 增加术语准确性检查
            if hasattr(self, 'terminology_manager'):
                terminology_score = self._check_terminology_accuracy(original, translation)
                score += terminology_score * 0.3
            
            # 2. 增加句式结构检查
            structure_score = self._check_sentence_structure(translation)
            score += structure_score * 0.2
            
            # 3. 增强数字和特殊字符检查
            special_content_score = self._check_special_content(original, translation)
            score += special_content_score * 0.2
            
            # 4. 增加语言流畅度检查
            fluency_score = self._check_language_fluency(translation)
            score += fluency_score * 0.2
            
            # 5. 增加上下文一致性检查
            context_score = self._check_context_consistency(translation)
            score += context_score * 0.1
            
            return min(1.0, score)
            
        except Exception as e:
            print(f"评分过程出错: {str(e)}")
            return 0.0

    def _check_terminology_accuracy(self, original: str, translation: str) -> float:
        """检查术语翻译准确性"""
        try:
            if not self.terminology_manager.terminology:
                return 1.0
            
            correct_terms = 0
            total_terms = 0
            
            for cn_term, en_term in self.terminology_manager.terminology.items():
                if cn_term in original:
                    total_terms += 1
                    if en_term in translation:
                        correct_terms += 1
                    
            return correct_terms / total_terms if total_terms > 0 else 1.0
        except:
            return 0.8

    def _is_valid_translation_with_reason(self, original_text: str, translated_text: str) -> Tuple[bool, str]:
        """带原因的翻译验证"""
        try:
            if not translated_text or not translated_text.strip():
                return False, "空翻译结果"
            
            # 移除 [START] 和 [END] 标签
            cleaned_text = re.sub(r'\[START\]|\[END\]', '', translated_text).strip()
            if not cleaned_text:
                return False, "清理后为空"
            
            # 检查是否与原文完全相同
            if cleaned_text == original_text:
                return False, "翻译结果与原文相同"
           
            if all(char in string.punctuation + ' ' for char in cleaned_text):
                return False, "仅包含标点符号"
            
            # 检查中文内容
            if self.prompt_manager.has_chinese(cleaned_text):
                return False, "包含中文字符"
            
            if len(cleaned_text.strip()) < 2:
                return False, "翻译过短"
            
            # 检查是否包含指令关键词
            instruction_keywords = ['must', 'requirement', 'output', 'translation:', 'example:']
            if any(keyword in cleaned_text.lower() for keyword in instruction_keywords):
                return False, "包含指令关键词"
            
            # 检查数字保护
            original_numbers = set(re.findall(r'-?\d+\.?\d*', original_text))
            translated_numbers = set(re.findall(r'-?\d+\.?\d*', cleaned_text))
            if not original_numbers.issubset(translated_numbers):
                return False, "数字未被正确保留"
            
            # 语言检测
            try:
                # 忽略短文本的语言检测
                if len(cleaned_text.split()) > 3: # 只对超过3个单词的文本进行语言检测
                    detected_lang = detect(cleaned_text)
                    if detected_lang != 'en':
                        return False, f"错误的语言类型: {detected_lang}"
            except:
                pass  # 忽略语言检测失败
            
            return True, "有效翻译"
            
        except Exception as e:
            return False, f"验证过程错误: {str(e)}"

    def _protect_special_content(self, text: str) -> Tuple[str, Dict[str, str]]:
        """增强的数字和特殊内容保护"""
        token_map = {}
        
        # 需要保护的模式
        patterns = [
            # 基本数字（整数、小数）
            r'-?\d+\.?\d*',
            # 分数
            r'\d+/\d+',
            # 科学计数法
            r'-?\d+\.?\d*[eE][+-]?\d+',
            # 带单位的数字
            r'-?\d+\.?\d*\s*[a-zA-Z]+',
            # 带括号的数字
            r'\(\d+\.?\d*\)',
            # 特殊格式（如 1-2, 1.2.3）
            r'\d+[-\.]\d+(?:[-\.]\d+)*'
        ]
        
        protected_text = text
        for pattern in patterns:
            def replace_match(match):
                token = f"__TOKEN_{len(token_map)}__"
                token_map[token] = match.group(0)
                return token
                
            protected_text = re.sub(pattern, replace_match, protected_text)
        
        return protected_text, token_map

    def _restore_special_content(self, text: str, token_map: Dict[str, str]) -> str:
        """恢复被保护的内容"""
        result = text
        # 按token长度降序排序，避免部分替换问题
        for token, value in sorted(token_map.items(), key=lambda x: len(x[0]), reverse=True):
            result = result.replace(token, value)
        return result

    def _adjust_translation_strategy(self, text: str, failure_reasons: List[str]) -> str:
        """根据失败原因调整翻译策略"""
        # 分析失败原因
        reason_counts = {}
        for reason in failure_reasons[-3:]:  # 只看最近的3次失败
            reason_counts[reason] = reason_counts.get(reason, 0) + 1
        
        most_common_reason = max(reason_counts.items(), key=lambda x: x[1])[0]
        
        # 根据不同的失败原因采取不同的策略
        if "包含中文字符" in most_common_reason:
            # 强制使用英文翻译策略
            return f"[FORCE_ENGLISH]{text}"
        elif "数字未被正确保留" in most_common_reason:
            # 增强数字保护
            return f"[PROTECT_NUMBERS]{text}"
        elif "仅包含标点符号" in most_common_reason:
            # 要求完整翻译
            return f"[FULL_TRANSLATION]{text}"
        else:
            # 默认策略：分段翻译
            return f"[SEGMENT]{text}"

    def _get_best_candidate(self, original_text: str, failure_reasons: List[str]) -> str:
        """在所有尝试中选择最佳结果"""
        try:
            # 最后一次尝试：直接分词翻译
            words = original_text.split()
            translated_parts = []
            
            for word in words:
                if not self._is_translatable(word):
                    translated_parts.append(word)
                    continue
                    
                # 使用最简单的翻译策略
                translation = self.translate_text_with_retry(word, False, 0)
                if self._is_valid_translation_with_reason(word, translation)[0]:
                    translated_parts.append(translation)
                else:
                    translated_parts.append(word)
            
            return ' '.join(translated_parts)
            
        except Exception as e:
            print(f"获取最佳候选失败: {str(e)}")
            return original_text

    def is_target_language(self, text: str) -> bool:
        """检查文本是否属于目标语种"""
        # 这里可以使用语言检测库，如langdetect，来检测文本的语言
        # 以下是一个简单的示例，假设目标语种为英文
        try:
            detected_language = detect(text)
            return detected_language == 'en'
        except:
            return False

    def batch_translate_texts(self, texts: List[str], use_terminology: bool = False) -> List[str]:
        """批量翻译文本（包含复检）"""
        futures = []
        results = [""] * len(texts)

        # 提交所有翻译任务到线程池
        for i, text in enumerate(texts):
            future = self.executor.submit(
                self.translate_text_with_verification, 
                text, 
                use_terminology
            )
            futures.append((i, future))

        # 按顺序获取结果
        for i, future in futures:
            try:
                results[i] = future.result()
            except Exception as e:
                print(f"翻译任务 {i} 失败: {str(e)}")
                results[i] = texts[i]

        return results

    def _clean_translation_output(self, text: str) -> str:
        """清理翻译输出，移除标签和不必要的内容"""
        try:
            if not text:
                return ""
            
            # 移除 [START] 和 [END] 标签
            cleaned = re.sub(r'\[START\]|\[END\]', '', text)
            
            # 移除多余的空白字符
            cleaned = ' '.join(cleaned.split())
            
            # 移除可能的指令关键词
            instruction_keywords = [
                'must', 'requirement', 'output:', 'translation:', 
                'example:', 'input:', 'strict requirements:'
            ]
            
            # 检查每一行是否包含指令关键词
            lines = cleaned.split('\n')
            valid_lines = []
            for line in lines:
                if not any(keyword in line.lower() for keyword in instruction_keywords):
                    valid_lines.append(line)
            
            cleaned = ' '.join(valid_lines)
            
            # 确保结果不为空
            if not cleaned.strip():
                return text
            
            return cleaned.strip()
            
        except Exception as e:
            print(f"清理翻译输出时出错: {str(e)}")
            return text

    def _detect_domain(self, text: str) -> str:
        """识别文本所属专业领域"""
        domains = {
            'technical': ['parameter', 'system', 'device', 'module'],
            'medical': ['treatment', 'patient', 'diagnosis', 'symptom'],
            'financial': ['revenue', 'profit', 'market', 'investment'],
            'legal': ['contract', 'agreement', 'clause', 'party']
        }
        
        domain_scores = {domain: 0 for domain in domains}
        words = text.lower().split()
        
        for word in words:
            for domain, keywords in domains.items():
                if word in keywords:
                    domain_scores[domain] += 1
                    
        if any(domain_scores.values()):
            return max(domain_scores.items(), key=lambda x: x[1])[0]
        return 'general'

    def _check_translation_consistency(self, translations: List[Dict[str, str]]) -> List[Dict[str, str]]:
        """检查并确保翻译的一致性"""
        try:
            # 建立术语映射表
            term_mapping = {}
            
            # 第一遍：收集所有翻译对应关系
            for item in translations:
                original = item['original']
                translation = item['translation']
                
                # 分词处理
                original_terms = self._extract_terms(original)
                translated_terms = self._extract_terms(translation)
                
                # 更新术语映射
                for o_term, t_term in zip(original_terms, translated_terms):
                    if o_term not in term_mapping:
                        term_mapping[o_term] = {}
                    if t_term not in term_mapping[o_term]:
                        term_mapping[o_term][t_term] = 0
                    term_mapping[o_term][t_term] += 1
            
            # 第二遍：统一翻译用语
            for item in translations:
                item['translation'] = self._unify_translation(
                    item['translation'],
                    term_mapping
                )
                
            return translations
            
        except Exception as e:
            print(f"一致性检查失败: {str(e)}")
            return translations

    def _auto_correct_translation(self, translation: str) -> str:
        """自动纠正常见翻译错误"""
        try:
            # 1. 修正标点符号使用
            translation = self._fix_punctuation(translation)
            
            # 2. 修正大小写使用
            translation = self._fix_capitalization(translation)
            
            # 3. 修正空格使用
            translation = self._fix_spacing(translation)
            
            # 4. 修正专业术语
            translation = self._fix_terminology(translation)
            
            return translation
            
        except Exception as e:
            print(f"自动纠错失败: {str(e)}")
            return translation

    def load_terminology_rules(self):
        """加载术语规则"""
        self.terminology_rules = {
            # 强制替换规则
            "force_replace": {
                "oven": "Puller",  # 统一使用 Puller
                "hot field": "Hot zone",  # 统一使用 Hot zone
            },
            # 上下文相关替换规则
            "context_replace": {
                "炉": {
                    "default": "Puller",
                    "contexts": {
                        "开炉": "open the Puller",
                        "合炉": "close the Puller",
                    }
                }
            }
        }

    def _process_special_markers(self, text: str) -> str:
        """处理文本中的特殊标记"""
        try:
            # 处理特殊标记
            markers = {
                '[START]': '',  # 移除开始标记
                '[END]': '',    # 移除结束标记
                '[FORCE_ENGLISH]': '',  # 移除强制英文标记
                '[PROTECT_NUMBERS]': '', # 移除数字保护标记
                '[SEGMENT]': ''  # 移除分段标记
            }
            
            # 移除所有特殊标记
            for marker, replacement in markers.items():
                text = text.replace(marker, replacement)
                
            # 清理多余的空格
            text = ' '.join(text.split())
            
            return text.strip()
            
        except Exception as e:
            print(f"处理特殊标记时出错: {str(e)}")
            return text

    def _get_translation(self, text: str, use_terminology: bool = False) -> str:
        """获取翻译结果"""
        try:
            # 如果文本为空，直接返回
            if not text.strip():
                return text
            
            # 构建翻译提示词
            messages = self.prompt_manager.get_translation_prompt(
                text=text,
                source_lang="Chinese",
                target_lang="English" if self.config.config["target_language"] == "英文" else "Chinese",
                context="ppt",
                is_verification=False
            )
            
            # 调用翻译客户端
            translation = self.client.translate(
                messages=[
                    {"role": "system", "content": messages["system"]},
                    {"role": "user", "content": messages["user"]}
                ],
                model=self.config.config["model"]
            )
            
            # 如果启用了术语库，应用术语规则
            if use_terminology and hasattr(self, 'terminology_manager'):
                translation = self.terminology_manager.apply_terminology(translation)
            
            # 清理翻译结果
            translation = self._clean_translation_output(translation)
            
            return translation
            
        except Exception as e:
            print(f"翻译失败: {str(e)}")
            return text

    def _restore_special_format(self, text: str) -> str:
        """恢复特殊格式"""
        try:
            # 恢复数字和单位
            text = re.sub(r'__NUM__(\d+)__UNIT__([a-zA-Z°℃]+)__END__', r'\1\2', text)
            
            # 恢复括号内容
            text = re.sub(r'__LEFT__([^_]+)__RIGHT__', r'（\1）', text)
            
            # 恢复其他可能的特殊格式
            special_formats = {
                '__DASH__': '-',
                '__COLON__': ':',
                '__SLASH__': '/',
                '__DOT__': '.',
                '__COMMA__': ','
            }
            
            for marker, symbol in special_formats.items():
                text = text.replace(marker, symbol)
            
            # 清理可能残留的保护标记
            text = re.sub(r'__[A-Z]+__', '', text)
            
            # 清理多余的空格
            text = ' '.join(text.split())
            
            return text.strip()
            
        except Exception as e:
            print(f"恢复特殊格式时出错: {str(e)}")
            return text

    def _validate_translation(self, translation: str, original: str) -> bool:
        """验证翻译结果的质量"""
        try:
            # 1. 基本验证
            if not translation or not translation.strip():
                return False
            
            # 2. 检查是否与原文完全相同
            if translation.strip() == original.strip():
                return False
            
            # 3. 检查目标语言
            if self.config.config["target_language"] == "英文":
                # 检查是否包含中文字符
                if self.prompt_manager.has_chinese(translation):
                    return False
                
                # 检查是否只包含标点符号
                if all(char in string.punctuation + ' ' for char in translation):
                    return False
                
                # 检查长度是否合理（不应过短）
                if len(translation.strip()) < 2:
                    return False
                
                # 检查是否包含指令关键词
                instruction_keywords = ['must', 'requirement', 'output', 'translation:', 'example:']
                if any(keyword in translation.lower() for keyword in instruction_keywords):
                    return False
            
            # 4. 检查数字保护
            original_numbers = set(re.findall(r'-?\d+\.?\d*', original))
            translated_numbers = set(re.findall(r'-?\d+\.?\d*', translation))
            if not original_numbers.issubset(translated_numbers):
                return False
            
            # 5. 检查术语准确性
            if not self._verify_terminology(translation):
                return False
            
            # 6. 检查语言流畅度
            if not self._verify_fluency(translation):
                return False
            
            # 7. 检查上下文一致性
            if not self._verify_context(translation):
                return False
            
            return True
            
        except Exception as e:
            print(f"翻译验证失败: {str(e)}")
            return False

    def _retry_translation(self, text: str) -> str:
        """当翻译验证失败时进行重试"""
        try:
            # 最多重试3次
            for i in range(3):
                # 使用不同的提示词变体
                messages = self.prompt_manager.get_variant_prompt(text, i)
                translation = self.client.translate(messages, self.config.config["model"])
                
                # 验证新的翻译结果
                if translation and self._validate_translation(translation, text):
                    return translation
                
                # 在重试之间添加短暂延迟
                time.sleep(1)
            
            # 如果所有重试都失败，返回原文
            return text
            
        except Exception as e:
            print(f"重试翻译失败: {str(e)}")
            return text
